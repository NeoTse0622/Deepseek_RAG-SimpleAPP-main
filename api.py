# coding=utf-8
from fastapi import FastAPI, HTTPException, UploadFile, File, Request
from fastapi.responses import JSONResponse
from pydantic import BaseModel
from typing import List, Dict
import os
import shutil
import datetime
import json  # 用于处理历史对话记录
import torch
from transformers import AutoModelForCausalLM, AutoTokenizer
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.vectorstores import FAISS
from langchain.chains import RetrievalQA
from langchain.document_loaders import TextLoader, PyPDFLoader, Docx2txtLoader, UnstructuredFileLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.prompts import PromptTemplate
from langchain.llms import HuggingFacePipeline
from transformers import pipeline
from pptx import Presentation  # 用于读取PPT文件

app = FastAPI()

# 全局变量 (根据实际情况修改)
MODEL_PATH = "./deepseek-llm-7b-chat"
DB_FAISS_PATH = "vectorstore/db_faiss"
UPLOAD_FOLDER = "uploads"  # 用于保存上传的知识库文件
HISTORY_FILE = "history.json"  # 用于保存历史对话记录
os.makedirs(UPLOAD_FOLDER, exist_ok=True)  # 确保上传文件夹存在

# 初始化模型、tokenizer 和向量数据库 (在应用启动时加载)
try:
    tokenizer = AutoTokenizer.from_pretrained(MODEL_PATH)
    model = AutoModelForCausalLM.from_pretrained(MODEL_PATH, device_map="auto", torch_dtype=torch.float16)

    pipe = pipeline("text-generation",
                    model=model,
                    tokenizer=tokenizer,
                    torch_dtype=torch.float16,
                    device_map="auto",
                    max_new_tokens=256,
                    do_sample=True,
                    top_p=0.9,
                    temperature=0.7,
                    num_return_sequences=1)

    llm = HuggingFacePipeline(pipeline=pipe)

    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2",
                                       model_kwargs={'device': 'cuda' if torch.cuda.is_available() else 'cpu'})
    db = FAISS.load_local(DB_FAISS_PATH, embeddings, allow_dangerous_deserialization=True)

    # 自定义 Prompt
    prompt_template = """使用以下上下文来回答最后的问题。如果你不知道答案，就说你不知道，不要试图编造答案。
    上下文：{context}
    问题：{question}
    有用的回答："""
    prompt = PromptTemplate(
        template=prompt_template, input_variables=["context", "question"]
    )

    qa_chain = RetrievalQA.from_chain_type(
        llm=llm,
        chain_type="stuff",
        retriever=db.as_retriever(search_kwargs={'k': 3}),
        return_source_documents=True,
        chain_type_kwargs={"prompt": prompt}
    )

    print("模型和向量数据库加载成功！")

except Exception as e:
    print(f"模型或向量数据库加载失败: {e}")
    raise  # 抛出异常，阻止应用启动

# 辅助函数
def load_document(file_path: str):
    """加载单个文档"""
    try:
        if file_path.endswith(".txt"):
            loader = TextLoader(file_path, encoding="utf-8")
        elif file_path.endswith(".pdf"):
            loader = PyPDFLoader(file_path)
        elif file_path.endswith(".docx") or file_path.endswith(".doc"):
            try:
                loader = Docx2txtLoader(file_path)
            except ImportError:
                print(f"docx2txt 未安装，尝试使用 UnstructuredFileLoader 加载 {file_path}")
                loader = UnstructuredFileLoader(file_path)
        elif file_path.endswith(".ppt") or file_path.endswith(".pptx"):
            # 使用 python-pptx 加载 PPT 文件
            text = ""
            try:
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text
            except Exception as e:
                print(f"加载 PPT 文件失败: {e}")
                raise
            # 将 PPT 内容视为一个文本文件加载
            loader = TextLoader(file_path, encoding="utf-8")
            loader.load()[0].page_content = text  # 将提取的文本内容写入 page_content
            return loader.load()  # 返回加载的文档列表
        else:
            loader = UnstructuredFileLoader(file_path)
        return loader.load()
    except Exception as e:
        print(f"加载文档 {file_path} 失败: {e}")
        raise

def create_vector_db_from_files(file_paths: List[str]):
    """从文件列表创建向量数据库"""
    documents = []
    for file_path in file_paths:
        try:
            documents.extend(load_document(file_path))
        except Exception as e:
            print(f"加载文档 {file_path} 失败: {e}")
            continue  # 忽略加载失败的文件

    if not documents:
        raise ValueError("没有成功加载任何文档")

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
    texts = text_splitter.split_documents(documents)

    if not texts:
        raise ValueError("没有成功分割任何文本块")

    global db  # 使用全局变量
    db = FAISS.from_documents(texts, embeddings)
    db.save_local(DB_FAISS_PATH)
    print("向量数据库创建/更新完成！")

def load_history():
    """加载历史对话记录"""
    try:
        with open(HISTORY_FILE, "r", encoding="utf-8") as f:
            try:
                history = json.load(f)
            except json.JSONDecodeError as e:
                print(f"加载历史记录失败，文件可能损坏: {e}")
                history = []  # 如果文件损坏，则创建一个空列表
    except FileNotFoundError:
        history = []
    return history

def save_history(history: List[Dict[str, str]]):
    """保存历史对话记录"""
    with open(HISTORY_FILE, "w", encoding="utf-8") as f:
        json.dump(history, f, ensure_ascii=False, indent=4)

# API 接口

@app.post("/uploadfiles/")
async def upload_files(files: List[UploadFile] = File(...)):
    """上传知识库文件"""
    file_paths = []
    for file in files:
        try:
            file_path = os.path.join(UPLOAD_FOLDER, file.filename)
            with open(file_path, "wb") as f:
                shutil.copyfileobj(file.file, f)
            file_paths.append(file_path)
        except Exception as e:
            return JSONResponse(status_code=500, content={"message": f"文件 {file.filename} 上传失败: {e}"})
        finally:
            file.file.close()

    try:
        create_vector_db_from_files(file_paths)
        return {"message": "知识库文件上传成功，向量数据库已更新！"}
    except ValueError as e:
        return JSONResponse(status_code=400, content={"message": str(e)})
    except Exception as e:
        return JSONResponse(status_code=500, content={"message": f"创建向量数据库失败: {e}"})

# 定义请求体模型
class ChatRequest(BaseModel):
    query: str

@app.post("/chat/")
async def chat(request: ChatRequest):
    """对话接口"""
    try:
        query = request.query  # 从请求体中获取 query
        response = await qa_chain.ainvoke({"query": query})  # 使用 ainvoke 进行异步调用
        print(f"qa_chain.ainvoke 的返回值: {response}")  # 打印返回值

        result = response["result"]

        # 提取 "有用的回答" 部分的内容
        useful_answer_prefix = "有用的回答："
        if useful_answer_prefix in result:
            useful_answer = result.split(useful_answer_prefix)[1].strip()  # 提取 "有用的回答" 之后的内容
        else:
            useful_answer = result  # 如果没有找到 "有用的回答"，则返回整个 result
        
        # 获取历史记录
        history = load_history()

        # 添加当前对话到历史记录
        history.append({"query": query, "response": useful_answer, "timestamp": datetime.datetime.now().isoformat()})

        # 保存更新后的历史记录
        save_history(history)

        # 只返回 "有用的回答" 部分的内容
        return JSONResponse(content={"response": useful_answer})
    except Exception as e:
        print(f"发生异常: {type(e).__name__} - {e}")  # 打印异常类型和消息
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/history/")
async def get_history():
    """获取历史对话记录"""
    history = load_history()
    return history

# 启动信息
@app.get("/")
async def read_root():
    return {"message": "RAG API is running!"}

@app.middleware("http")
async def log_requests(request: Request, call_next):
    """记录请求和响应的中间件"""
    print(f"请求 URL: {request.url}")
    print(f"请求方法: {request.method}")
    print(f"请求头: {request.headers}")
    try:
        body = await request.body()
        print(f"请求体: {body.decode()}")
    except Exception as e:
        print(f"无法读取请求体: {e}")

    response = await call_next(request)

    print(f"响应状态码: {response.status_code}")
    print(f"响应头: {response.headers}")
    try:
        response_body = b""
        async for chunk in response.body_iterator():
            response_body += chunk
        print(f"响应体: {response_body.decode()}")
    except Exception as e:
        print(f"无法读取响应体: {e}")

    return response

